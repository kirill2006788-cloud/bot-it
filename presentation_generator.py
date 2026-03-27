#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Генератор презентаций для Telegram бота
Создает PowerPoint презентации с помощью Claude AI и генерации изображений
"""

import io
import logging
from typing import Dict, List, Optional, Tuple
import json
import re

# PowerPoint
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Word
try:
    from docx import Document
    WORD_AVAILABLE = True
except ImportError:
    WORD_AVAILABLE = False

# PDF
try:
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib.colors import HexColor
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

logger = logging.getLogger(__name__)

class PresentationGenerator:
    """Класс для генерации презентаций с помощью AI"""
    
    def __init__(self, claude_helper=None, image_generator=None):
        self.claude_helper = claude_helper
        self.image_generator = image_generator
        
        # Шаблоны для разных типов презентаций
        self.presentation_templates = {
            'business': {
                'title': 'Бизнес презентация',
                'slides_count': 8,
                'structure': [
                    'Титульный слайд',
                    'Проблема/Возможность',
                    'Решение',
                    'Рынок',
                    'Бизнес-модель',
                    'Конкуренты',
                    'Финансы',
                    'Следующие шаги'
                ]
            },
            'educational': {
                'title': 'Образовательная презентация',
                'slides_count': 10,
                'structure': [
                    'Титульный слайд',
                    'Введение',
                    'Основные понятия',
                    'Теория',
                    'Примеры',
                    'Практическое применение',
                    'Преимущества',
                    'Недостатки',
                    'Заключение',
                    'Вопросы'
                ]
            },
            'technical': {
                'title': 'Техническая презентация',
                'slides_count': 12,
                'structure': [
                    'Титульный слайд',
                    'Обзор',
                    'Архитектура',
                    'Технологии',
                    'Компоненты',
                    'Интеграция',
                    'Производительность',
                    'Безопасность',
                    'Тестирование',
                    'Развертывание',
                    'Мониторинг',
                    'Заключение'
                ]
            },
            'marketing': {
                'title': 'Маркетинговая презентация',
                'slides_count': 9,
                'structure': [
                    'Титульный слайд',
                    'Целевая аудитория',
                    'Проблема клиента',
                    'Наше решение',
                    'Преимущества',
                    'Кейсы',
                    'Цены',
                    'Команда',
                    'Контакты'
                ]
            },
            'research': {
                'title': 'Исследовательская презентация',
                'slides_count': 11,
                'structure': [
                    'Титульный слайд',
                    'Цель исследования',
                    'Методология',
                    'Гипотезы',
                    'Данные',
                    'Анализ',
                    'Результаты',
                    'Выводы',
                    'Ограничения',
                    'Рекомендации',
                    'Спасибо'
                ]
            }
        }
    
    def analyze_text_with_claude(self, text: str, presentation_type: str = 'business') -> Dict:
        """Анализ текста с помощью Claude AI для создания структуры презентации"""
        if not self.claude_helper:
            return self._create_default_structure(text, presentation_type)
        
        try:
            template = self.presentation_templates.get(presentation_type, self.presentation_templates['business'])
            
            prompt = f"""Проанализируй следующий текст и создай структуру презентации типа "{template['title']}".

Текст: {text}

Создай презентацию из {template['slides_count']} слайдов со следующей структурой:
{', '.join(template['structure'])}

Для каждого слайда предоставь:
1. Заголовок слайда
2. Основной текст (2-3 предложения)
3. Ключевые моменты (3-5 пунктов)
4. Описание изображения/диаграммы (если нужно)

Ответь в формате JSON:
{{
    "title": "Название презентации",
    "slides": [
        {{
            "slide_number": 1,
            "title": "Заголовок слайда",
            "content": "Основной текст",
            "bullet_points": ["Пункт 1", "Пункт 2", "Пункт 3"],
            "image_description": "Описание изображения",
            "needs_diagram": false
        }}
    ]
}}"""

            response = self.claude_helper.ask_question(prompt)
            
            # Пытаемся извлечь JSON из ответа
            json_match = re.search(r'\{.*\}', response, re.DOTALL)
            if json_match:
                json_str = json_match.group()
                return json.loads(json_str)
            else:
                # Если не удалось извлечь JSON, создаем структуру по умолчанию
                return self._create_default_structure(text, presentation_type)
                
        except Exception as e:
            logger.error(f"Ошибка анализа с Claude: {e}")
            return self._create_default_structure(text, presentation_type)
    
    def _create_default_structure(self, text: str, presentation_type: str) -> Dict:
        """Создание структуры презентации по умолчанию"""
        template = self.presentation_templates.get(presentation_type, self.presentation_templates['business'])
        
        # Разбиваем текст на части
        sentences = text.split('. ')
        slides_count = min(template['slides_count'], len(template['structure']))
        
        slides = []
        for i in range(slides_count):
            slide_text = sentences[i * len(sentences) // slides_count:(i + 1) * len(sentences) // slides_count]
            content = '. '.join(slide_text)
            
            slides.append({
                'slide_number': i + 1,
                'title': template['structure'][i],
                'content': content[:200] + '...' if len(content) > 200 else content,
                'bullet_points': [
                    f"Ключевой момент {j + 1}" 
                    for j in range(3)
                ],
                'image_description': f"Иллюстрация для {template['structure'][i].lower()}",
                'needs_diagram': i in [2, 4, 6]  # Некоторые слайды нуждаются в диаграммах
            })
        
        return {
            'title': f"Презентация: {template['title']}",
            'slides': slides
        }
    
    def generate_presentation(self, text: str, presentation_type: str = 'business', user_id: int = 0) -> Tuple[bytes, str]:
        """Генерация полной презентации"""
        if not PPTX_AVAILABLE:
            raise Exception("PowerPoint обработка недоступна. Установите python-pptx")
        
        try:
            # Анализируем текст с помощью Claude
            structure = self.analyze_text_with_claude(text, presentation_type)
            
            # Создаем презентацию
            prs = Presentation()
            
            # Удаляем стандартный слайд
            slide_layout = prs.slide_layouts[0]
            prs.slides._sldIdLst.clear()
            
            # Создаем слайды
            for slide_data in structure['slides']:
                slide = prs.slides.add_slide(slide_layout)
                
                # Заголовок
                title_shape = slide.shapes.title
                title_shape.text = slide_data['title']
                
                # Содержимое
                content_shape = slide.placeholders[1]
                text_frame = content_shape.text_frame
                text_frame.clear()
                
                # Основной текст
                p = text_frame.paragraphs[0]
                p.text = slide_data['content']
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(64, 64, 64)
                
                # Маркированный список
                for bullet_point in slide_data['bullet_points']:
                    p = text_frame.add_paragraph()
                    p.text = bullet_point
                    p.font.size = Pt(12)
                    p.font.color.rgb = RGBColor(96, 96, 96)
                    p.level = 0
                
                # Добавляем изображение если нужно
                if slide_data.get('needs_diagram', False) and self.image_generator:
                    try:
                        # Генерируем изображение для слайда
                        image_prompt = f"Professional diagram for presentation slide: {slide_data['image_description']}"
                        image_buffer, _, _ = self.image_generator.generate_with_settings(
                            image_prompt, user_id, enhance_with_claude=False
                        )
                        
                        # Добавляем изображение на слайд
                        image_buffer.seek(0)
                        left = Inches(6)
                        top = Inches(2)
                        width = Inches(3)
                        height = Inches(2)
                        
                        slide.shapes.add_picture(image_buffer, left, top, width, height)
                    except Exception as e:
                        logger.warning(f"Не удалось добавить изображение: {e}")
            
            # Сохраняем презентацию
            presentation_buffer = io.BytesIO()
            prs.save(presentation_buffer)
            presentation_buffer.seek(0)
            
            filename = f"{structure['title'].replace(' ', '_')}.pptx"
            
            return presentation_buffer.getvalue(), filename
            
        except Exception as e:
            raise Exception(f"Ошибка генерации презентации: {e}")
    
    def create_html_presentation(self, text: str, presentation_type: str = 'business') -> str:
        """Создание HTML версии презентации"""
        try:
            structure = self.analyze_text_with_claude(text, presentation_type)
            
            html_content = f"""<!DOCTYPE html>
<html lang="ru">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{structure['title']}</title>
    <style>
        body {{
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            margin: 0;
            padding: 20px;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            min-height: 100vh;
        }}
        .container {{
            max-width: 1200px;
            margin: 0 auto;
            background: white;
            border-radius: 15px;
            box-shadow: 0 20px 40px rgba(0,0,0,0.1);
            overflow: hidden;
        }}
        .header {{
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: white;
            padding: 30px;
            text-align: center;
        }}
        .header h1 {{
            margin: 0;
            font-size: 2.5em;
            font-weight: 300;
        }}
        .slide {{
            padding: 40px;
            border-bottom: 1px solid #eee;
            page-break-after: always;
        }}
        .slide:last-child {{
            border-bottom: none;
        }}
        .slide h2 {{
            color: #333;
            font-size: 2em;
            margin-bottom: 20px;
            border-bottom: 3px solid #667eea;
            padding-bottom: 10px;
        }}
        .slide p {{
            font-size: 1.2em;
            line-height: 1.6;
            color: #666;
            margin-bottom: 20px;
        }}
        .bullet-points {{
            list-style: none;
            padding: 0;
        }}
        .bullet-points li {{
            padding: 10px 0;
            border-left: 4px solid #667eea;
            padding-left: 20px;
            margin-bottom: 10px;
            background: #f8f9ff;
            border-radius: 0 8px 8px 0;
        }}
        .slide-number {{
            position: absolute;
            top: 20px;
            right: 20px;
            background: #667eea;
            color: white;
            padding: 10px 15px;
            border-radius: 20px;
            font-weight: bold;
        }}
        @media print {{
            body {{ background: white; }}
            .container {{ box-shadow: none; }}
        }}
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <h1>{structure['title']}</h1>
        </div>
"""
            
            for slide_data in structure['slides']:
                html_content += f"""
        <div class="slide">
            <div class="slide-number">{slide_data['slide_number']}</div>
            <h2>{slide_data['title']}</h2>
            <p>{slide_data['content']}</p>
            <ul class="bullet-points">
"""
                for bullet_point in slide_data['bullet_points']:
                    html_content += f"                <li>{bullet_point}</li>\n"
                
                html_content += """            </ul>
        </div>
"""
            
            html_content += """
    </div>
</body>
</html>"""
            
            return html_content
            
        except Exception as e:
            raise Exception(f"Ошибка создания HTML презентации: {e}")
    
    def create_word_presentation(self, text: str, presentation_type: str = 'business') -> Tuple[bytes, str]:
        """Создание Word версии презентации"""
        try:
            structure = self.analyze_text_with_claude(text, presentation_type)
            
            # Создаем новый документ Word
            doc = Document()
            
            # Настройка стилей
            title_style = doc.styles['Title']
            title_style.font.size = Pt(24)
            title_style.font.bold = True
            
            heading_style = doc.styles['Heading 1']
            heading_style.font.size = Pt(18)
            heading_style.font.bold = True
            
            # Титульная страница
            doc.add_heading(structure['title'], 0)
            doc.add_paragraph(f"Тип презентации: {self.presentation_templates[presentation_type]['title']}")
            doc.add_paragraph(f"Количество слайдов: {len(structure['slides'])}")
            doc.add_page_break()
            
            # Создаем слайды
            for slide_data in structure['slides']:
                # Заголовок слайда
                doc.add_heading(f"Слайд {slide_data['slide_number']}: {slide_data['title']}", level=1)
                
                # Основной текст
                if slide_data['content']:
                    doc.add_paragraph(slide_data['content'])
                
                # Маркированный список
                if slide_data['bullet_points']:
                    for bullet_point in slide_data['bullet_points']:
                        doc.add_paragraph(bullet_point, style='List Bullet')
                
                # Разрыв страницы между слайдами
                if slide_data['slide_number'] < len(structure['slides']):
                    doc.add_page_break()
            
            # Сохраняем документ
            doc_buffer = io.BytesIO()
            doc.save(doc_buffer)
            doc_buffer.seek(0)
            
            filename = f"{structure['title'].replace(' ', '_')}.docx"
            
            return doc_buffer.getvalue(), filename
            
        except Exception as e:
            raise Exception(f"Ошибка создания Word презентации: {e}")
    
    def create_pdf_presentation(self, text: str, presentation_type: str = 'business') -> Tuple[bytes, str]:
        """Создание PDF версии презентации"""
        try:
            structure = self.analyze_text_with_claude(text, presentation_type)
            
            # Создаем PDF документ
            pdf_buffer = io.BytesIO()
            doc = SimpleDocTemplate(pdf_buffer, pagesize=A4)
            
            # Стили
            styles = getSampleStyleSheet()
            
            # Кастомные стили
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Title'],
                fontSize=24,
                spaceAfter=30,
                alignment=1,  # Center
                textColor=HexColor('#2E86AB')
            )
            
            slide_title_style = ParagraphStyle(
                'SlideTitle',
                parent=styles['Heading1'],
                fontSize=18,
                spaceAfter=20,
                textColor=HexColor('#2E86AB')
            )
            
            content_style = ParagraphStyle(
                'Content',
                parent=styles['Normal'],
                fontSize=12,
                spaceAfter=15,
                leftIndent=20
            )
            
            bullet_style = ParagraphStyle(
                'Bullet',
                parent=styles['Normal'],
                fontSize=11,
                spaceAfter=8,
                leftIndent=30,
                bulletIndent=20
            )
            
            # Содержимое документа
            story = []
            
            # Титульная страница
            story.append(Paragraph(structure['title'], title_style))
            story.append(Spacer(1, 20))
            story.append(Paragraph(f"<b>Тип презентации:</b> {self.presentation_templates[presentation_type]['title']}", content_style))
            story.append(Paragraph(f"<b>Количество слайдов:</b> {len(structure['slides'])}", content_style))
            story.append(PageBreak())
            
            # Создаем слайды
            for slide_data in structure['slides']:
                # Заголовок слайда
                slide_title = f"Слайд {slide_data['slide_number']}: {slide_data['title']}"
                story.append(Paragraph(slide_title, slide_title_style))
                
                # Основной текст
                if slide_data['content']:
                    story.append(Paragraph(slide_data['content'], content_style))
                
                # Маркированный список
                if slide_data['bullet_points']:
                    for bullet_point in slide_data['bullet_points']:
                        story.append(Paragraph(f"• {bullet_point}", bullet_style))
                
                # Разрыв страницы между слайдами
                if slide_data['slide_number'] < len(structure['slides']):
                    story.append(PageBreak())
            
            # Строим PDF
            doc.build(story)
            pdf_buffer.seek(0)
            
            filename = f"{structure['title'].replace(' ', '_')}.pdf"
            
            return pdf_buffer.getvalue(), filename
            
        except Exception as e:
            raise Exception(f"Ошибка создания PDF презентации: {e}")
    
    def get_presentation_types(self) -> Dict[str, str]:
        """Получение доступных типов презентаций"""
        return {
            key: template['title'] 
            for key, template in self.presentation_templates.items()
        }
    
    def format_presentation_info(self) -> str:
        """Форматирование информации о типах презентаций"""
        types = self.get_presentation_types()
        
        text = "🎯 **Генератор презентаций**\n\n"
        text += "Создавайте профессиональные презентации с помощью AI!\n\n"
        text += "**Доступные типы:**\n\n"
        
        for key, title in types.items():
            template = self.presentation_templates[key]
            text += f"📊 **{title}**\n"
            text += f"   • {template['slides_count']} слайдов\n"
            text += f"   • {', '.join(template['structure'][:3])}...\n\n"
        
        text += "**Как использовать:**\n"
        text += "1. Выберите тип презентации\n"
        text += "2. Опишите тему или вставьте текст\n"
        text += "3. Получите готовую презентацию!\n\n"
        text += "**Что получите:**\n"
        text += "• PowerPoint файл (.pptx)\n"
        text += "• HTML версию для просмотра\n"
        text += "• Структурированный контент\n"
        text += "• Изображения и диаграммы\n"
        
        return text
